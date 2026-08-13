from fastapi import FastAPI, File, UploadFile, Form
from fastapi.responses import JSONResponse
import magic
import os
import shutil
import zipfile
import tempfile
import uvicorn

app = FastAPI(title="NexusAI File Processor")

MAX_TEXT_CHARS = 200_000

SUPPORTED_TYPES = {
    "text/plain": "text",
    "text/csv": "text",
    "text/html": "text",
    "text/markdown": "text",
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/msword": "docx",
    "application/vnd.ms-excel": "xlsx",
    "application/vnd.ms-powerpoint": "pptx",
    "image/jpeg": "image",
    "image/png": "image",
    "image/gif": "image",
    "image/webp": "image",
    "image/bmp": "image",
    "audio/mpeg": "audio",
    "audio/wav": "audio",
    "audio/ogg": "audio",
    "audio/mp4": "audio",
    "video/mp4": "video",
    "video/avi": "video",
    "video/quicktime": "video",
    "video/webm": "video",
    "application/zip": "zip",
}

# Plain-text files read directly (by extension — Windows browsers often report
# these as application/octet-stream, so magic alone is unreliable).
TEXT_EXTENSIONS = {
    ".txt", ".csv", ".tsv", ".md", ".markdown", ".html", ".htm", ".json",
    ".xml", ".yml", ".yaml", ".log", ".ini", ".cfg", ".conf", ".toml", ".env",
    ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".c", ".h", ".cpp", ".hpp",
    ".cs", ".go", ".rb", ".php", ".sh", ".bash", ".sql", ".css", ".scss",
    ".less", ".ipynb", ".swift", ".kt", ".lua", ".pl", ".r", ".dart", ".vue",
    ".svelte", ".rst", ".tex", ".gitignore", ".dockerfile",
}

TEMP_DIR = "/app/temp"
os.makedirs(TEMP_DIR, exist_ok=True)


def _read_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read(MAX_TEXT_CHARS)


def _extract_pdf(path: str) -> str:
    try:
        import fitz
        doc = fitz.open(path)
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text
    except ImportError:
        return "PDF processing requires PyMuPDF"
    except Exception as e:
        return f"PDF error: {e}"


def _iter_block_items(parent):
    """Yield Paragraph and Table objects in document order from a python-docx
    parent (Document, _Cell, or header/footer). This is what python-docx does
    internally to build .paragraphs / .tables, but it also works for cells so
    table content is not skipped."""
    from docx.document import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    if isinstance(parent, Document):
        parent_elm = parent.element.body
    else:
        parent_elm = getattr(parent, "_tc", None)
        if parent_elm is None:
            parent_elm = getattr(parent, "_element", None)
    if parent_elm is None:
        return
    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def _walk_docx_content(parent):
    """Recursively yield every paragraph, including those inside tables and
    nested tables (merged cells de-duplicated). Body order is preserved.

    Merged cells are de-duplicated by the XML element itself — NOT by id(),
    which is unsafe here: when the walk is consumed without retaining the
    paragraphs (e.g. a list comprehension), the per-cell proxy objects can be
    garbage-collected and their memory addresses reused, which makes real
    cells look like duplicates and silently drops their text."""
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for block in _iter_block_items(parent):
        if isinstance(block, Paragraph):
            yield block
        else:
            seen_cells = set()
            for row in block.rows:
                for cell in row.cells:
                    tc = cell._tc
                    if tc in seen_cells:
                        continue  # merged cell repeats in later rows
                    seen_cells.add(tc)
                    yield from _walk_docx_content(cell)


def _extract_docx(path: str) -> str:
    """Extract text from a .docx covering body paragraphs, tables (the content
    of many Word documents lives in tables), headers/footers and text boxes.
    The plain .paragraphs read misses all of those."""
    try:
        from docx import Document
        import re
        doc = Document(path)

        parts: list[str] = []
        for p in _walk_docx_content(doc):
            if p.text and p.text.strip():
                parts.append(p.text)

        # Headers and footers often carry page titles / metadata.
        for section in doc.sections:
            for hf in (section.header, section.footer):
                try:
                    if hf.is_linked_to_previous:
                        continue
                except Exception:
                    pass
                try:
                    for p in _walk_docx_content(hf):
                        if p.text and p.text.strip():
                            parts.append(p.text)
                except Exception:
                    pass

        # Text boxes (floating/inline) are not exposed by python-docx — pull
        # them straight from the document XML.
        xml = doc.element.xml
        for box in re.findall(r"<w:txbxContent>(.*?)</w:txbxContent>", xml, re.S):
            text = "".join(re.findall(r"<w:t[^>]*>([^<]*)</w:t>", box)).strip()
            if text:
                parts.append(text)

        return "\n".join(parts)
    except ImportError:
        return "DOCX processing requires python-docx"
    except Exception as e:
        return f"DOCX error: {e}"


def _rtf_to_text(path: str) -> str:
    """Minimal RTF → plain text. Many ".doc" files are actually RTF (Word's
    default when saving from older tools). Control words are dropped;
    \\par/\\line/\\row become newlines and \\cell/\\tab become tabs, so table
    layouts survive as tab-separated rows. Metadata destinations (\\fonttbl,
    \\colortbl, \\stylesheet, \\info, and anything marked \\*) are skipped."""
    try:
        with open(path, "r", encoding="cp1252", errors="replace") as f:
            rtf = f.read(MAX_TEXT_CHARS)
    except Exception as e:
        return f"RTF read error: {e}"

    SKIP_DESTINATIONS = {
        "fonttbl", "colortbl", "stylesheet", "info", "datastore", "filetbl",
        "listtable", "listoverridetable", "revtbl", "rsidtbl", "themedata",
        "colorschememapping", "latentstyles", "pgptbl", "generator",
        "object", "pict", "nonshppict", "oleclsid",
    }

    def _group_is_skipped(rtf: str, start: int, n: int) -> bool:
        """After '{' at `start`, does this group open a destination to skip?"""
        j = start + 1
        while j < n and rtf[j] in " \r\n\t":
            j += 1
        if j >= n or rtf[j] != "\\":
            return False
        k = j + 1
        if k < n and rtf[k] == "*":
            return True  # {\\*...} destinations are explicitly ignorable
        while k < n and rtf[k].isalpha():
            k += 1
        return rtf[j + 1 : k] in SKIP_DESTINATIONS

    out: list[str] = []
    skip_depth = 0
    i, n = 0, len(rtf)
    while i < n:
        c = rtf[i]
        if c == "\\":
            i += 1
            if i >= n:
                break
            c = rtf[i]
            if c in "\\{}":
                if skip_depth == 0:
                    out.append(c)
                i += 1
                continue
            if c == "'":
                if skip_depth == 0:
                    try:
                        out.append(chr(int(rtf[i + 1 : i + 3], 16)))
                    except ValueError:
                        pass
                i += 3
                continue
            if not c.isalpha():
                i += 1
                continue
            j = i
            while j < n and rtf[j].isalpha():
                j += 1
            word = rtf[i:j]
            # signed numeric parameter (e.g. \\u-1234)
            k = j
            neg = 1
            if k < n and rtf[k] == "-":
                neg = -1
                k += 1
            while k < n and rtf[k].isdigit():
                k += 1
            try:
                param = neg * int(rtf[j:k] or "0")
            except ValueError:
                param = 0
            if k < n and rtf[k] == " ":
                k += 1  # space terminates a control word
            i = k
            if skip_depth:
                continue
            if word == "u" and param:
                out.append(chr(param & 0xFFFF))
                # skip the fallback char used for non-Unicode readers (\\'3f)
                if rtf[i : i + 2] == "\\'" and i + 4 <= n:
                    i += 4
            elif word in ("par", "line", "row", "sect"):
                out.append("\n")
            elif word in ("cell", "tab"):
                out.append("\t")
            elif word == "bullet":
                out.append("\u2022")
            elif word == "emdash":
                out.append("\u2014")
            elif word == "endash":
                out.append("\u2013")
            elif word in ("ldblquote", "lquote"):
                out.append("\u201c" if word == "ldblquote" else "\u2018")
            elif word in ("rdblquote", "rquote"):
                out.append("\u201d" if word == "rdblquote" else "\u2019")
            continue
        if c == "{":
            if _group_is_skipped(rtf, i, n):
                skip_depth += 1
            i += 1
            continue
        if c == "}":
            if skip_depth:
                skip_depth -= 1
            i += 1
            continue
        if c in "\r\n":
            i += 1
            continue
        if skip_depth == 0:
            out.append(c)
        i += 1
    return "".join(out).strip()


def _doc_binary_text(path: str) -> str:
    """Extract text from a genuine binary .doc (OLE compound file, Word 97+).

    python-docx only understands the modern OOXML format, so this parses the
    WordDocument stream directly: the CLX piece table (in the 0Table/1Table
    stream) maps each character range to a byte range in the WordDocument
    stream, encoded as UTF-16LE or ANSI (compressed) pieces. Table cell text
    lives inline in the same text stream, so it is included — cell marks
    become tabs and paragraph marks become newlines."""
    import struct

    try:
        import olefile
    except ImportError:
        return "Legacy .doc requires olefile (pip install olefile)"

    try:
        ole = olefile.OleFileIO(path)
    except Exception as e:
        return f"Legacy .doc parse error: {e}"

    try:
        if not ole.exists("WordDocument"):
            return "[No WordDocument stream in .doc]"
        word_doc = ole.openstream("WordDocument").read()

        n_fib = struct.unpack_from("<H", word_doc, 0x02)[0]
        if 0 < n_fib < 105:
            return "[Word 95-or-earlier .doc layout not supported]"

        # fWhichTblStm (FIB flags word, bit 0x0200) selects 1Table vs 0Table.
        flags = struct.unpack_from("<H", word_doc, 0x0A)[0]
        table_name = "1Table" if (flags & 0x0200) else "0Table"
        if not ole.exists(table_name):
            table_name = "0Table" if table_name == "1Table" else "1Table"
        if not ole.exists(table_name):
            return "[No table stream in .doc]"
        table = ole.openstream(table_name).read()

        # fcClx/lcbClx live in FibRgFcLcb97 at FIB offset 0x01A2 (Word 97+).
        fc_clx, lcb_clx = struct.unpack_from("<ii", word_doc, 0x01A2)
        if fc_clx < 0 or lcb_clx <= 0:
            return "[No piece table (CLX) in .doc]"
        clx = table[fc_clx : fc_clx + lcb_clx]

        # CLX = zero or more Prc entries (0x01) then one Pcdt (0x02).
        pos = 0
        pcdt = None
        while pos < len(clx):
            tag = clx[pos]
            if tag == 0x01:  # Prc: cbGrpprl (2B) + grpprl
                pos += 3 + struct.unpack_from("<H", clx, pos + 1)[0]
            elif tag == 0x02:  # Pcdt: lcb (4B) + PlcPcd
                lcb = struct.unpack_from("<I", clx, pos + 1)[0]
                pcdt = clx[pos + 5 : pos + 5 + lcb]
                break
            else:
                break
        if not pcdt:
            return "[Unsupported .doc piece table]"

        # PlcPcd: (n+1) CPs (4B each) followed by n PCDs (8B each).
        n = (len(pcdt) - 4) // 12
        cps = [struct.unpack_from("<I", pcdt, 4 * idx)[0] for idx in range(n + 1)]
        pcd_base = 4 * (n + 1)

        chunks: list[str] = []
        for idx in range(n):
            fc = struct.unpack_from("<I", pcdt, pcd_base + 8 * idx + 2)[0]
            compressed = bool(fc & 0x40000000)
            fc &= 0x3FFFFFFF
            length = cps[idx + 1] - cps[idx]  # characters in the piece
            if compressed:
                # ANSI piece: 1 byte per char; fc is a byte offset divided by 2.
                raw = word_doc[fc // 2 : fc // 2 + length]
                chunks.append(raw.decode("cp1252", errors="replace"))
            else:
                # UTF-16LE piece: 2 bytes per char.
                raw = word_doc[fc : fc + 2 * length]
                chunks.append(raw.decode("utf-16-le", errors="replace"))

        text = "".join(chunks)
        # Paragraph (\r), cell (\x07), line (\x0b) and page (\x0c) breaks.
        text = text.replace("\r", "\n").replace("\x07", "\t")
        text = text.replace("\x0b", "\n").replace("\x0c", "\n")
        # Word fields are \x13 instruction \x14 result \x15 — keep only the
        # result (e.g. HYPERLINK shows the URL/label, not the field code).
        clean: list[str] = []
        in_field = False
        for ch in text:
            if ch == "\x13":
                in_field = True
            elif ch in ("\x14", "\x15"):
                in_field = False
            elif not in_field:
                clean.append(ch)
        text = "".join(clean)
        # Footnote/endnote reference marks (\x03/\x04), separators and NULs.
        for ch in ("\x01", "\x02", "\x03", "\x04", "\x00"):
            text = text.replace(ch, "")
        return text.strip()
    finally:
        ole.close()


def _extract_doc(path: str) -> str:
    """Legacy binary .doc (Word 97-2003) — python-docx cannot open these.
    Sniff the container: RTF-with-a-.doc-extension is parsed as RTF, genuine
    OLE compound files are parsed via the WordDocument piece table."""
    try:
        with open(path, "rb") as f:
            head = f.read(4)
    except Exception as e:
        return f"Legacy .doc read error: {e}"
    if head[:2] == b"{\"":
        return _rtf_to_text(path)
    if head == b"\xd0\xcf\x11\xe0":
        return _doc_binary_text(path)
    return "[Unsupported legacy .doc format]"


def _extract_xlsx(path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None]
                if vals:
                    rows.append("\t".join(vals))
            parts.append(f"--- Sheet: {ws.title} ---\n" + "\n".join(rows))
        wb.close()
        return "\n".join(parts)
    except ImportError:
        return "XLSX processing requires openpyxl"
    except Exception as e:
        return f"XLSX error: {e}"


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs)
                        if t:
                            texts.append(t)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        texts.append("\t".join(cell.text for cell in row.cells))
            parts.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        return "\n".join(parts)
    except ImportError:
        return "PPTX processing requires python-pptx"
    except Exception as e:
        return f"PPTX error: {e}"


def _extract_image(path: str) -> str:
    try:
        import easyocr
        reader = easyocr.Reader(["en"])
        result = reader.readtext(path)
        return "\n".join(r[1] for r in result)
    except ImportError:
        return "Image OCR requires EasyOCR (not installed in this image)"
    except Exception as e:
        return f"Image OCR error: {e}"


def _extract_audio(path: str) -> str:
    try:
        import whisper
        model = whisper.load_model("base")
        result = model.transcribe(path)
        return result.get("text", "")
    except ImportError:
        return "Audio transcription requires openai-whisper (not installed in this image)"
    except Exception as e:
        return f"Audio transcription error: {e}"


def extract_zip(zip_path: str, depth: int = 0) -> str:
    """Recursively extract text from every readable file inside a ZIP."""
    if depth > 3:
        return "[Nested archives too deep]"
    parts = []
    try:
        with zipfile.ZipFile(zip_path) as zf:
            infos = [i for i in zf.infolist() if not i.is_dir()]
            infos.sort(key=lambda i: i.filename)
            for info in infos:
                name = info.filename
                base = os.path.basename(name)
                if not base or base.startswith(".") or base == ".DS_Store" or "__MACOSX" in name:
                    continue
                sub_dir = tempfile.mkdtemp(dir=TEMP_DIR)
                try:
                    target = os.path.join(sub_dir, base)
                    with zf.open(name) as src, open(target, "wb") as dst:
                        shutil.copyfileobj(src, dst)
                    if os.path.splitext(base)[1].lower() == ".zip":
                        text = extract_zip(target, depth + 1)
                    else:
                        text = extract_text_from_file(target, depth + 1)
                    if text and text.strip():
                        parts.append(f"\n===== {name} =====\n{text}")
                except Exception as e:
                    parts.append(f"\n===== {name} =====\n[Error: {e}]")
                finally:
                    shutil.rmtree(sub_dir, ignore_errors=True)
    except zipfile.BadZipFile:
        return "[Not a valid ZIP archive]"
    except Exception as e:
        return f"[ZIP error: {e}]"

    out = "\n".join(parts)
    return out[:MAX_TEXT_CHARS]


def extract_text_from_file(path: str, depth: int = 0) -> str:
    """Extract readable text from a file, dispatching on extension first (more
    reliable than MIME for files uploaded from Windows browsers), then MIME."""
    ext = os.path.splitext(path)[1].lower()
    try:
        mime = magic.from_file(path, mime=True)
    except Exception:
        mime = ""

    if ext in TEXT_EXTENSIONS or mime.startswith("text/"):
        return _read_text(path)
    if ext == ".pdf" or mime == "application/pdf":
        return _extract_pdf(path)
    if ext == ".docx" or "wordprocessingml" in mime:
        return _extract_docx(path)
    if ext == ".doc" or mime == "application/msword":
        return _extract_doc(path)
    if ext in (".xlsx", ".xls") or "spreadsheetml" in mime or mime == "application/vnd.ms-excel":
        return _extract_xlsx(path)
    if ext in (".pptx", ".ppt") or "presentationml" in mime or mime == "application/vnd.ms-powerpoint":
        return _extract_pptx(path)
    if ext == ".zip" or mime == "application/zip":
        return extract_zip(path, depth)
    if mime.startswith("image/"):
        return _extract_image(path)
    if mime.startswith("audio/"):
        return _extract_audio(path)
    return "[No extractable text for this file type]"


@app.get("/health")
async def health():
    return {"status": "healthy"}


@app.get("/supported-types")
async def supported_types():
    return {"types": list(SUPPORTED_TYPES.keys())}


@app.post("/process")
async def process_file(
    file: UploadFile = File(...),
    generate_embeddings: bool = Form(False),
    include_metadata: bool = Form(False)
):
    file_type = magic.from_buffer(await file.read(2048), mime=True)
    await file.seek(0)

    category = SUPPORTED_TYPES.get(file_type, "unknown")

    temp_path = os.path.join(TEMP_DIR, os.path.basename(file.filename or "upload"))
    with open(temp_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    try:
        extracted_text = extract_text_from_file(temp_path)
    except Exception as e:
        extracted_text = f"Error processing file: {str(e)}"
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

    response = {"text": extracted_text[:MAX_TEXT_CHARS]}

    if include_metadata:
        response["metadata"] = {
            "filename": file.filename,
            "mime_type": file_type,
            "category": category,
            "size": file.size if hasattr(file, "size") else 0,
        }

    return JSONResponse(content=response)


@app.post("/process-path")
async def process_file_path(
    file_path: str,
    generate_embeddings: bool = False,
    include_metadata: bool = False
):
    if not os.path.exists(file_path):
        return JSONResponse(content={"error": "File not found"}, status_code=404)

    try:
        extracted_text = extract_text_from_file(file_path)
    except Exception as e:
        extracted_text = f"Error processing file: {str(e)}"

    return JSONResponse(content={"text": extracted_text[:MAX_TEXT_CHARS]})


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
